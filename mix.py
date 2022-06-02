from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st
from io import BytesIO
import plotly.express as px
import pandas as pd
import numpy as np
import os
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor

# emojis: https://www.webfx.com/tools/emoji-cheat-sheet/
st.set_page_config(page_title="Technician Rating Dashboard", page_icon=":bar_chart:", layout="wide")
st.header('Latar Belakang')
st.markdown("<div style='text-align: justify;'>Kepuasan pelanggan dapat diukur salah satunya dari aspek kepuasan pelayanan terhadap pelanggan. Jika pelayanan yang diterima oleh pelanggan tidak/memenuhi harapan mereka, maka pelanggan akan memberikan feedback sesuai apa yang mereka peroleh. Data rating teknisi diperlukan untuk mengevaluasi teknisi apakah mereka telah menjalankan tugas dengan baik atau tidak.</div>", unsafe_allow_html=True)
st.markdown('''
Alur untuk mengisi survey Data Teknisi:
* Isi FAQ.
* Melaporkan gangguan.
* Menjadwalkan kapan untuk perbaikan perangkat.
* Repair Request (Riwayat Tracking Teknisi berangkat hingga problem solved).
* Setelah problem solved, pelanggan mengisi survey Data Teknisi (rating, alasan memberikan rating tersebut, regional, witel, id teknisi, data diri pengguna, dan lain-lain).
''')

#uploader file
uploaded = st.file_uploader("Upload file", type=['csv','xlsx'])
st.write(uploaded)

global df
if uploaded:
    # Check MIME type of the uploaded file
    if uploaded.type == "text/csv":
        df = pd.read_csv(uploaded)
    else:
        df = pd.read_excel(uploaded)
try:
    df['responses.createdAt'] = pd.to_datetime(df['responses.createdAt'])
    df['responses_date'] = df['responses.createdAt'].dt.strftime('%B %Y')
    df.replace({'responses.witel' : {'JATENG UTARA  (SEMARANG)': 'SEMARANG', 'BABEL': 'BANGKA BELITUNG', 'RIDAR': 'RIAU DARATAN','KALTENG':'KALIMANTAN TENGAH',
    'JATENG SELATAN  (MAGELANG)':'MAGELANG', 'SUMUT': 'SUMATERA UTARA', 'KALBAR': 'KALIMANTAN BARAT', 'RIKEP': 'RIAU KEPULAUAN','SUMATERA SELATAN (PALEMBANG)':'PALEMBANG',
    'SUMSEL': 'SUMATERA SELATAN', 'SERANG': 'BANTEN', 'BANDUNGBRT': 'BANDUNG BARAT', 'JATENG BARAT SELATAN  (PWKT)': 'PURWOKERTO','SUMATERA BARAT (PADANG)':'PADANG',
    'KALTARA': 'KALIMANTAN UTARA', 'JATIM TENGAH  (KEDIRI)': 'KEDIRI', 'SUMBAR': 'SUMATERA BARAT', 'KALSEL': 'KALIMANTAN SELATAN','JABAR TENGAH (BANDUNG BRT)': 'BANDUNG BARAT',
    'JATENG TIMUR SELATAN  (SOLO)': 'SOLO', 'JATIM SELATAN  (MALANG)': 'MALANG', 'JATIM TENGAH TIMUR  (SIDOARJO)': 'SIDOARJO','KALTENG  (PALANGKARAYA)':'PALANGKARAYA',
    'BALI SELATAN  (DENPASAR)': 'DENPASAR', 'JAKSEL': 'JAKARTA SELATAN', 'JAKTIM': 'JAKARTA TIMUR', 'JABAR BARAT UTARA (BEKASI)': 'BEKASI','KALSEL  (BANJARMASIN)':'BANJARMASIN',
    'SULTENG': 'SULAWESI TENGAH', 'SUMUT BARAT (MEDAN)': 'MEDAN', 'BALI UTARA  (SINGARAJA)': 'SINGARAJA', 'DI YOGYAKARTA': 'YOGYAKARTA','BANTEN TIMUR (TANGERANG)': 'TANGERANG',
    'SULTRA': 'SULAWESI TENGGARA', 'NTT  (KUPANG)': 'KUPANG', 'JATIMSEL TIMUR  (PASURUAN)': 'PASURUAN', 'JABAR BARAT (BOGOR)': 'BOGOR','JAKUT':'JAKARTA UTARA',
    'DATEL MAKASAR':'MAKASAR','JAKBAR': 'JAKARTA BARAT','JATENG BARAT UTARA  (PKLG)':'PEKALONGAN','JATENG TIMUR UTARA  (KUDUS)':'KUDUS',
    'JATIM BARAT  (MADIUN)':'MADIUN','KALTIMSEL  (BALIKPAPAN)':'BALIKPAPAN','NTB  (MATARAM)':'MATARAM'}}, inplace=True)
except Exception as e:
    print(e)

#preprocessing
col_dict = {'responses.region': 'responses_region', 'responses.witel':'responses_witel', 'responses.rate':'responses_rate'}
df.columns = [col_dict.get(x, x) for x in df.columns]
#df.drop(df.columns[25:42],axis=1,inplace=True)

# ---- READ EXCEL ----
#data 1
data = pd.read_excel('technician rating.xlsx')
st.title('Data Rating Teknisi')

# ---- SIDEBAR ----
st.sidebar.header("Please Filter Here:")
date = st.sidebar.multiselect(
    "Select the Date:",
    options=df["responses_date"].unique(),
    default=df["responses_date"].unique()
)

region = st.sidebar.multiselect(
    "Select the Region:",
    options=df["responses_region"].unique(),
    default=df["responses_region"].unique()
)

witel_ = st.sidebar.multiselect(
    "Select the Witel:",
    options=df["responses_witel"].unique(),
    default=df["responses_witel"].unique()
)

# use selected values from widgets to filter dataset down to only the rows we need
df = df.query(
    "responses_date == @date & responses_region == @region & responses_witel==@witel_")
df = df.astype(str)
st.dataframe(df)

st.header('Dataset yang telah dipilih')
st.write('Data Dimension: ' + str(df.shape[0]) + ' rows and ' + str(df.shape[1]) + ' columns.')

# ---- MAINPAGE ----
st.title(":bar_chart: Technician Rating Dashboard")
st.markdown("##")
# simple description
st.write('In this dashboard we will analyze the Technician Rating data from NPS. '
           'These charts are interactive')

# TOP KPI
total_review = int(df["responses_rate"].count())
average_rating = round(data["responses.rate"].mean(), 2)
star_rating = ":star:" * int(round(average_rating, 0))

left_column, right_column = st.columns(2)
with left_column:
    st.subheader("Total Review:")
    st.subheader(f" {total_review:,}")
with right_column:
    st.subheader("Average Rating:")
    st.subheader(f"{average_rating} {star_rating}")

st.markdown("""---""")

#############################
#date
jumlah_date = df['responses_date'].value_counts().rename_axis('date').reset_index(name='count').sort_values(by=["date"])
ratee = df['responses_rate'].value_counts()
date = pd.concat([df['responses_date'],ratee], axis=1)
date = date.groupby(['responses_date'], as_index=False).sum()
date.rename(columns = {'responses_date': 'date'}, inplace = True)
date = date.join(jumlah_date.set_index('date'), on='date').sort_values(by=['count'], ascending=False)

#region
jumlah_rg = df['responses_region'].value_counts().rename_axis('region').reset_index(name='count').sort_values(by=["region"])
regi = pd.concat([df['responses_region'],ratee], axis=1)
regi = regi.groupby(['responses_region'], as_index=False).sum()
regi.rename(columns = {'responses_region': 'region'}, inplace = True)
regi = regi.join(jumlah_rg.set_index('region'), on='region').sort_values(by=['count'], ascending=False)

#responseCount
jumlah_res = df['responseCount'].value_counts().rename_axis('responseCount').reset_index(name='count').sort_values(by=['responseCount'])
res = pd.concat([df['responseCount'],ratee], axis=1)
res = res.groupby(['responseCount'], as_index=False).sum()
res = res.join(jumlah_res.set_index('responseCount'), on='responseCount').sort_values(by=['count'], ascending=False)

#incidence count
jumlah_inc = df['incidentCount'].value_counts().rename_axis('incidentCount').reset_index(name='count').sort_values(by=['incidentCount'])
inc = pd.concat([df['incidentCount'],ratee], axis=1)
inc = inc.groupby(['incidentCount'], as_index=False).sum()
#inc.rename(columns = {'responses_region': 'region'}, inplace = True)
inc = inc.join(jumlah_inc.set_index('incidentCount'), on='incidentCount').sort_values(by=['count'], ascending=False)

#pie
jumlah_rt = df['responses_rate'].value_counts().rename_axis('rate').reset_index(name='count').sort_values(by=["rate"])
#rt = pd.concat([df['responses_rate'],ratee], axis=1)
rt = df.groupby(['responses_rate'], as_index=False).sum()
rt.rename(columns = {'responses_rate': 'rate'}, inplace = True)
rt = rt.join(jumlah_rt.set_index('rate'), on='rate').sort_values(by=['count'], ascending=False)

#witel
banyak_witel = df['responses_witel'].value_counts().rename_axis('witel').reset_index(name='count').sort_values(by=['witel'])
Witel_raw = pd.concat([df['responses_witel'], ratee], axis=1)
Witel = Witel_raw.groupby(['responses_witel'], as_index=False).sum()
Witel.rename(columns = {'responses_witel': 'witel'}, inplace = True)
Witel = Witel.join(banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)

#sto
banyak_sto = df['responses.sto'].value_counts().rename_axis('sto').reset_index(name='count').sort_values(by=['sto'])
STO_raw = pd.concat([df['responses.sto'], ratee], axis=1)
STO = STO_raw.groupby(['responses.sto'], as_index=False).sum()
STO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
STO = STO.join(banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)

# Rating Extraction
rating = pd.get_dummies(df['responses_rate'], columns=['responses_rate'], prefix='rating', prefix_sep=' ')

# RateXRegion
RateXRegion_raw = pd.concat([df['responses_region'], rating], axis=1)
RateXRegion = RateXRegion_raw.groupby(['responses_region'], as_index=False).sum()
RateXRegion.rename(columns = {'responses_region': 'region'}, inplace = True)
# RateXWitel
RateXWitel_raw = pd.concat([df['responses_witel'], rating], axis=1)
RateXWitel = RateXWitel_raw.groupby(['responses_witel'], as_index=False).sum()
RateXWitel.rename(columns = {'responses_witel': 'witel'}, inplace = True)
RateXWitel_fix = RateXWitel.join(banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
# RateXSTO
RateXSTO_raw = pd.concat([df['responses.sto'], rating], axis=1)
RateXSTO = RateXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
RateXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
RateXSTO = RateXSTO.join(banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)
RateXSTO = RateXSTO.astype({
                    'rating 1':int,
                    'rating 2':int,
                    'rating 3':int,
                    'rating 4':int,
                    'rating 5':int,
})
# Feedback extraction
b = df['responses.selectedOptions']
#feedback = ','.join( str(a) for a in b)
replace = {"Arrive on time": "Datang tepat waktu", "Friendly":"Ramah", "Explain the cause of the disruption":"Menjelaskan penyebab gangguan", "The technician's appearance / uniform is presentable" : "Penampilan/ seragam teknisi rapi",
"Resolve the disruption quickly": "Menyelesaikan gangguan dengan cepat","Slow resolution of disruption":"Penyelesaikan gangguan lambat"}
replacer = replace.get
abc = [replacer(n, n) for n in b]
abc = (pd.DataFrame(b.tolist())
                      .fillna('')
                      .astype(str)
                      .agg(','.join, 1)
                      .str.strip(','))
feedback = abc.str.get_dummies(sep=',')
#replace({'responses.selectedOptions' : {'Arrive on time': 'Datang tepat waktu', "Friendly":"Ramah", "Explain the cause of the disruption":
#"Menjelaskan penyebab gangguan", "The technician's appearance / uniform is presentable" : "Penampilan/ seragam teknisi rapi",
#"Resolve the disruption quickly": "Menyelesaikan gangguan dengan cepat","Slow resolution of disruption":"Penyelesaikan gangguan lambat"}}, inplace=True)

#feedback = [w.replace('Arrive on time','Datang tepat waktu') for w in feedback]
#feedback = [w.replace('Friendly','Ramah') for w in feedback]
#feedback = [w.replace("Explain the cause of the disruption","Menjelaskan penyebab gangguan") for w in feedback]
#feedback = [w.replace("The technician's appearance / uniform is presentable", "Penampilan/ seragam teknisi rapi") for w in feedback]
#feedback = [w.replace("Resolve the disruption quickly", "Menyelesaikan gangguan dengan cepat") for w in feedback]
#feedback = [w.replace("Slow resolution of disruption","Penyelesaikan gangguan lambat") for w in feedback]

feedback.replace('Arrive on time','Datang tepat waktu')
feedback.replace('Friendly','Ramah') 
feedback.replace("Explain the cause of the disruption","Menjelaskan penyebab gangguan")
feedback.replace("The technician's appearance / uniform is presentable", "Penampilan/ seragam teknisi rapi")
feedback.replace("Resolve the disruption quickly", "Menyelesaikan gangguan dengan cepat")
feedback.replace("Slow resolution of disruption","Penyelesaikan gangguan lambat")

ordered_feedback = ['Menyelesaikan gangguan dengan cepat',
                    'Datang tepat waktu',
                    'Ramah',
                    'Menjelaskan penyebab gangguan',
                    'Penampilan/ seragam teknisi rapi',
                    'Penyelesaikan gangguan lambat',
                    'Tidak datang tepat waktu',
                    'Kasar dan tidak sopan',
                    'Tidak berseragam/ tidak rapi',
                    'Menjelaskan layanan dengan baik',
                    'Menyelesaikan pemasangan dengan cepat']
#feedback = feedback[ordered_feedback]

# FeedbackXRegion
FeedbackXRegion_raw = pd.concat([df['responses_region'], feedback], axis=1)
FeedbackXRegion = FeedbackXRegion_raw.groupby(['responses_region'], as_index=False).sum()
# FeedbackXWitel
FeedbackXWitel_raw = pd.concat([df['responses_witel'], feedback], axis=1)
FeedbackXWitel = FeedbackXWitel_raw.groupby(['responses_witel'], as_index=False).sum()
FeedbackXWitel.rename(columns = {'responses_witel': 'witel'}, inplace = True)
FeedbackXWitel_fix = FeedbackXWitel.join(banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
# FeedbackXSTO
FeedbackXSTO_raw = pd.concat([df['responses.sto'], feedback], axis=1)
FeedbackXSTO = FeedbackXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
FeedbackXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
FeedbackXSTO = FeedbackXSTO.join(banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)

# extract date
bulan = pd.concat([df['responses_date'], rating], axis=1)
bulan = bulan.groupby(['responses_date'], as_index=False).sum()
bulan.rename(columns = {'responses_date': 'date'}, inplace = True)

# function to call graph
def plotlygraph(str_data, df, str_mode):
    if str_data=='head':
        fig = px.bar(df.head(10),
                     x=df.head(10).columns[0],
                     y=df.head(10).columns[1:-1],
                     barmode=str_mode
                     #title="Wide-Form Input",
                    )
    elif str_data=='tail':
        fig = px.bar(df.tail(10),
                     x=df.tail(10).columns[0],
                     y=df.tail(10).columns[1:-1],
                     barmode=str_mode
                    )
    else:
        fig = px.bar(df,
                     x=df.columns[0],
                     y=df.columns[1:],
                     barmode=str_mode
                    )
    st.plotly_chart(fig, use_container_width=True)   
    return

####################################################################################################################################################
st.header('Ringkasan Data Teknisi')
#1. review_by_date
review_by_date = (
    df.groupby(by=["responses_date"]).count()[["responses_rate"]].sort_values(by="responses_rate")
)
fig_date = px.bar(
    review_by_date,
    x="responses_rate",
    y=review_by_date.index,
    orientation="h",
    title="<b>Review By Date</b>",
    color_discrete_sequence=["#0083B8"] * len(review_by_date),
    template="plotly_white",text_auto=True
)
fig_date.update_layout(
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis_title="Count", 
    yaxis_title="Date")
st.plotly_chart(fig_date)

#2. Average Rate By Date
fig_avg_rd = px.histogram(x=df["responses_date"], y=df["responses_rate"],histfunc="avg",
title="<b>Average Rate By Date </b>",text_auto=True)
fig_avg_rd.update_layout(xaxis_title="Date", yaxis_title="Average Rate")
st.plotly_chart(fig_avg_rd)

#3. Review by Region
review_by_region = (
    df.groupby(by=["responses_region"]).count()[["responses_rate"]].sort_values(by="responses_rate")
)
fig_region = px.bar(
    review_by_region,
    x="responses_rate",
    y=review_by_region.index,
    orientation="h",
    title="<b>Review By Region</b>",
    color_discrete_sequence=["#0083B8"] * len(review_by_region),
    template="plotly_white",text_auto=True
)
fig_region.update_layout(
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis_title="Count", 
    yaxis_title="Region")
st.plotly_chart(fig_region)

#4. Average Rate By Region
fig_avg_rg = px.histogram(x=df["responses_region"], y=df["responses_rate"],histfunc="avg",
title="<b>Average Rate By Regional </b>",text_auto=True)
fig_avg_rg.update_layout(xaxis_title="Region", yaxis_title="Average Rate")
st.plotly_chart(fig_avg_rg)

#5. Rating Teknisi
#Rating Teknisi
rate = df['responses_rate'].value_counts()
labels = df['responses_rate'].unique()
fig_rate = px.pie(df,values=rate,names=labels,
    title="<b>Rating Teknisi </b>")
fig_rate.update_traces(hoverinfo='label+percent', textinfo='value')
st.plotly_chart(fig_rate)

#6. Many per Ticket Reopen
ticket = df['totalTicketReopenCount'].value_counts()
label1 = df['totalTicketReopenCount'].unique()
fig_avg_rate_tr = px.pie(df,values=ticket,names=label1,
    title="<b>Banyaknya Pelanggan Melakukan Ticket Reopen </b>")
fig_avg_rate_tr.update_traces(hoverinfo='label+percent', textinfo='value')
st.plotly_chart(fig_avg_rate_tr)

#7. Many Reopen Count
reopen = df['responses.reopenCount'].value_counts()
label2 = df['responses.reopenCount'].unique()
fig_avg_rate_r = px.pie(df,values=reopen,names=label2,
    title="<b>Banyaknya pelanggan melakukan Reopen </b>")
fig_avg_rate_r.update_traces(hoverinfo='label+percent', textinfo='value')
st.plotly_chart(fig_avg_rate_r)

#8. Response Count
fig_responseCount = px.histogram(x=df['responseCount'].unique(), y=df['responseCount'].value_counts(), 
title='<b>Response Count</b>',text_auto=True)
fig_responseCount.update_layout(xaxis_title="Response", yaxis_title="Count")
st.plotly_chart(fig_responseCount)

#9. Incidence Count
fig_incCount = px.line(x=df['incidentCount'].unique(), y=df['incidentCount'].value_counts()
, title='<b>Incident Count</b>')
fig_incCount.update_layout(xaxis_title="Incidence", yaxis_title="Count")
st.plotly_chart(fig_incCount)

#10. TOP 10 WITEL
witel_order = df['responses_witel'].value_counts().head(10)
fig_wtl = px.bar(
    witel_order,
    x='responses_witel',
    y=witel_order.index,orientation="h",
    title="<b>Top 10 Witel</b>",
    color_discrete_sequence=["#0083B8"] * len(witel_order),
    template="plotly_white",text_auto=True
)
fig_wtl.update_layout(
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis_title="Count", 
    yaxis_title="Witel")
st.plotly_chart(fig_wtl)

#11. TOP 10 STO
sto_order = df['responses.sto'].value_counts().head(10)
fig_sto = px.bar(
    sto_order,
    x='responses.sto',
    y=sto_order.index,orientation="h",
    title="<b>Top 10 STO</b>",
    color_discrete_sequence=["#0083B8"] * len(sto_order),
    template="plotly_white",text_auto=True
)
fig_sto.update_layout(
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis_title="Count", 
    yaxis_title="STO")
st.plotly_chart(fig_sto)

#####
st.header('Ringkasan Rating Teknisi Per Region')
plotlygraph(None, RateXRegion, 'group')

st.header('Ringkasan Rating Teknisi Per Witel')
st.subheader('Top 10 Rating')
plotlygraph('head', RateXWitel_fix, 'group')
st.subheader('Bottom 10 Rating')
plotlygraph('tail', RateXWitel_fix, 'group')

st.header('Ringkasan Rating Teknisi Per STO')
# st.write(RateXSTO)
st.subheader('Top 10 Rating')
plotlygraph('head', RateXSTO, 'group')
st.subheader('Bottom 10 Rating')
plotlygraph('tail', RateXSTO, 'group')

####
st.header('Ringkasan Feedback Teknisi Per Region')
plotlygraph(None, FeedbackXRegion, 'group')

st.header('Ringkasan Feedback Teknisi Per Witel')
st.subheader('Top 10 Feedback')
plotlygraph('head', FeedbackXWitel_fix, 'group')
st.subheader('Bottom 10 Feedback')
plotlygraph('tail', FeedbackXWitel_fix, 'group')

st.header('Ringkasan Feedback Teknisi Per STO')
# st.write(FeedbackXSTO)
st.subheader('Top 10 Feedback')
plotlygraph('head', FeedbackXSTO, 'group')
st.subheader('Bottom 10 Feedback')
plotlygraph('tail', FeedbackXSTO, 'group')

st.header('Ringkasan Rating Teknisi Per Bulan')
plotlygraph(None, bulan, 'group')
####################################################################################################################################################

# MAKE PPTX FILE
prs = Presentation('template.pptx')
# prs.slide_width = Inches(16)
# prs.slide_height = Inches(9)
slide_reg = prs.slide_layouts[11]

# Make bar clustered slide
def bar(df, str_title):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    title.text = str_title
    # title1.text_frame.paragraphs[0].font.name = "Arial"
    chart_data = ChartData()
    chart_data.categories = list(df.iloc[:,0])
    for i in range(1,len(df.columns)):
        chart_data.add_series(df.columns[i], (df.iloc[:,i]))
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'

    points =chart.plots[0].series[0].points
    fill = points[0].format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,255,255)
    return

# Make pie slide
def pie(df, str_title):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    title.text = str_title
    # title1.text_frame.paragraphs[0].font.name = "Arial"
    chart_data = ChartData()
    chart_data.categories = list(df.iloc[:,0])
    for i in range(1,len(df.columns)):
        chart_data.add_series(df.columns[i], (df.iloc[:,i]))
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'
    return

# Make fbXreg slide
def fbreg(df, str_title):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    title.text = str_title
    # title1.text_frame.paragraphs[0].font.name = "Arial"
    chart_data = ChartData()
    chart_data.categories = list(df.iloc[:,0])
    for i in range(1,len(df.columns)):
        chart_data.add_series(df.columns[i], (df.iloc[:,i]))
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = RGBColor(0,0,0)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'
    return

def addSeries(head, df, str_title):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    # title3.top = Inches(3)
    title.text = str_title
    chart_data = ChartData()
    if head:
        chart_data.categories = list(df.head(10).iloc[:,0])
        for i in range (1,len(df.columns)):
            chart_data.add_series(df.head(10).columns[i], (df.head(10).iloc[:,i]))
    else:
        chart_data.categories = list(df.tail(10).iloc[:,0])
        for i in range (1,len(df.columns)):
            chart_data.add_series(df.tail(10).columns[i], (df.tail(10).iloc[:,i]))

    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'
    return

#plotlygraph(None,date,'group')
bar(date,'Ringkasan Rating Teknisi Per Periode')
bar(regi,'Banyaknya Review Rating Teknisi Per Regional')
fbreg(res,'Response Count')
fbreg(inc,'Incidence Count')
addSeries(True, Witel , "Top 10 Review Ringkasan Rating Teknisi Per Witel")
addSeries(True, STO, "Top 10 Review Ringkasan Rating Teknisi Per STO")
fbreg(RateXRegion, 'Ringkasan Rating Teknisi Per Region')
addSeries(True, RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Top 10")
addSeries(False, RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Bottom 10")
addSeries(True, RateXSTO, "Ringkasan Rating Teknisi Per STO Top 10")
addSeries(False, RateXSTO, "Ringkasan Rating Teknisi Per STO Bottom 10")
fbreg(FeedbackXRegion, 'Ringkasan Feedback Teknisi Per Region')
addSeries(True, FeedbackXWitel_fix, "Ringkasan Feedback Teknisi Per Witel Top 10")
addSeries(False, FeedbackXWitel_fix, "Ringkasan Feedback Teknisi Per Witel Bottom 10")
addSeries(True, FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Top 10")
addSeries(False, FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Bottom 10")
fbreg(bulan,'Ringkasan Rating Teknisi Per Bulan')

last_slide = prs.slides.add_slide(prs.slide_layouts[13])

# save the output into binary form
binary_output = BytesIO()
prs.save(binary_output) 
st.download_button(label = 'Download ppt',
                   data = binary_output.getvalue(),
                   file_name = 'Technician Rate.pptx')

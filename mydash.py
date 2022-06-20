from doctest import DocFileTest
import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor

st.set_page_config(page_title="Technician Rating Dashboard", page_icon=":bar_chart:", layout="wide")

uploaded = st.sidebar.file_uploader("Upload File", type=['csv','xlsx'])
global df
if uploaded:
    # Check MIME type of the uploaded file
    if uploaded.type == "text/csv":
        df = pd.read_csv(uploaded)
    else:
        df = pd.read_excel(uploaded)


st.title("Technician Rating")  # add a title
st.header('Latar Belakang')
st.markdown('<div style="text-align: justify"> Kepuasan pelanggan dapat diukur salah satunya dari aspek kepuasan pelayanan terhadap pelanggan. Jika pelayanan yang diterima oleh pelanggan tidak/memenuhi harapan mereka, maka pelanggan akan memberikan feedback sesuai apa yang mereka peroleh. Data rating teknisi diperlukan untuk mengevaluasi teknisi apakah mereka telah menjalankan tugas dengan baik atau tidak. </div>', unsafe_allow_html=True)
st.markdown('''
Alur untuk mengisi survey Data Teknisi:
* Isi FAQ.
* Melaporkan gangguan.
* Menjadwalkan kapan untuk perbaikan perangkat.
* Repair Request (Riwayat Tracking Teknisi berangkat hingga problem solved).
* Setelah problem solved, pelanggan mengisi survey Data Teknisi (rating, alasan memberikan rating tersebut, regional, witel, id teknisi, data diri pengguna, dan lain-lain).
''')


#@st.cache
# df = pd.read_csv('technicianrate.csv', sep=';')
# df.drop(df.columns[32:],axis=1,inplace=True)
# cols = [8,21,22]
# df.drop(df.columns[cols],axis=1,inplace=True)
date = pd.to_datetime(df['responses.createdAt']).dt.strftime("%Y-%m")
df['indiHomeNum'] = df['indiHomeNum'].astype(str)
df['responses.region'] = 'region ' + df['indiHomeNum'].str[1]
# known_reg = ['1','2','3','4','5','6','7']
# for df['responses.indiHomeNum'].str[1] in known_reg:
#     df['responses.region'] = 'region ' + df['responses.indiHomeNum'].str[1]
# else:
#     df['responses.region'] = 'Undetected region ' + df['responses.indiHomeNum'].str[1]
# st.write(df[['responses.indiHomeNum', 'responses.region']])


witel_mapper = {'JATENG UTARA  (SEMARANG)': 'SEMARANG', 'BABEL': 'BANGKA BELITUNG', 'RIDAR': 'RIAU DARATAN','KALTENG':'KALIMANTAN TENGAH',
    'JATENG SELATAN  (MAGELANG)':'MAGELANG', 'SUMUT': 'SUMATERA UTARA', 'KALBAR': 'KALIMANTAN BARAT', 'RIKEP': 'RIAU KEPULAUAN','SUMATERA SELATAN (PALEMBANG)':'PALEMBANG',
    'SUMSEL': 'SUMATERA SELATAN', 'SERANG': 'BANTEN', 'BANDUNGBRT': 'BANDUNG BARAT', 'JATENG BARAT SELATAN  (PWKT)': 'PURWOKERTO','SUMATERA BARAT (PADANG)':'PADANG',
    'KALTARA': 'KALIMANTAN UTARA', 'JATIM TENGAH  (KEDIRI)': 'KEDIRI', 'SUMBAR': 'SUMATERA BARAT', 'KALSEL': 'KALIMANTAN SELATAN','JABAR TENGAH (BANDUNG BRT)': 'BANDUNG BARAT',
    'JATENG TIMUR SELATAN  (SOLO)': 'SOLO', 'JATIM SELATAN  (MALANG)': 'MALANG', 'JATIM TENGAH TIMUR  (SIDOARJO)': 'SIDOARJO','KALTENG  (PALANGKARAYA)':'PALANGKARAYA',
    'BALI SELATAN  (DENPASAR)': 'DENPASAR', 'JAKSEL': 'JAKARTA SELATAN', 'JAKTIM': 'JAKARTA TIMUR', 'JABAR BARAT UTARA (BEKASI)': 'BEKASI','KALSEL  (BANJARMASIN)':'BANJARMASIN',
    'SULTENG': 'SULAWESI TENGAH', 'SUMUT BARAT (MEDAN)': 'MEDAN', 'BALI UTARA  (SINGARAJA)': 'SINGARAJA', 'DI YOGYAKARTA': 'YOGYAKARTA','BANTEN TIMUR (TANGERANG)': 'TANGERANG',
    'SULTRA': 'SULAWESI TENGGARA', 'NTT  (KUPANG)': 'KUPANG', 'JATIMSEL TIMUR  (PASURUAN)': 'PASURUAN', 'JABAR BARAT (BOGOR)': 'BOGOR','JAKUT':'JAKARTA UTARA',
    'DATEL MAKASAR':'MAKASAR','JAKBAR': 'JAKARTA BARAT','JATENG BARAT UTARA  (PKLG)':'PEKALONGAN','JATENG TIMUR UTARA  (KUDUS)':'KUDUS',
    'JATIM BARAT  (MADIUN)':'MADIUN','KALTIMSEL  (BALIKPAPAN)':'BALIKPAPAN','NTB  (MATARAM)':'MATARAM'}

df['responses.witel'] = df['responses.witel'].map(witel_mapper).fillna(df['responses.witel'])
witel = df['responses.witel'].unique()
# selected_witel = st.sidebar.multiselect('Witel', witel, witel[:1])

# st.sidebar.header('Input Feature')
# df_selected_feature = df['responses.witel'].isin(selected_witel)
# df_selected_feature = df[(df.columns.isin(selected_cols)) & (df['responses.witel_new'].unique().isin(selected_witel))]

# st.header('Dataset yang telah dipilih')
# st.write('Data Dimension: ' + str(df[df_selected_feature].shape[0]) + ' rows and ' + str(df[df_selected_feature].shape[1]) + ' columns.')
# st.dataframe(df[df_selected_feature])

banyak_witel = df['responses.witel'].value_counts().rename_axis('witel').reset_index(name='count').sort_values(by=['witel'])
banyak_sto = df['responses.sto'].value_counts().rename_axis('sto').reset_index(name='count').sort_values(by=['sto'])
banyak_region = df['responses.region'].value_counts().rename_axis('region').reset_index(name='count').sort_values(by=['region'])

# Rating Extraction
rating = pd.get_dummies(df['responses.rate'], columns=['responses.rate'], prefix='rating', prefix_sep=' ')

# average per tiket reopen
avgr_ticketreop = pd.concat([df['responses.reopenCount'], rating], axis=1)
sum_avgr_ticketreop = avgr_ticketreop.groupby(['responses.reopenCount'], as_index=False).sum()
mean_avgr_ticketreop = avgr_ticketreop.groupby(['responses.reopenCount'], as_index=False).mean()
# RateXRegion
RateXRegion_raw = pd.concat([df['responses.region'], rating], axis=1)
RateXRegion = RateXRegion_raw.groupby(['responses.region'], as_index=False).sum()
# RateXWitel
RateXWitel_raw = pd.concat([df['responses.witel'], rating], axis=1)
RateXWitel = RateXWitel_raw.groupby(['responses.witel'], as_index=False).sum()
RateXWitel.rename(columns = {'responses.witel': 'witel'}, inplace = True)
RateXWitel_fix = RateXWitel.join(banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
# RateXSTO
RateXSTO_raw = pd.concat([df['responses.sto'], rating], axis=1)
RateXSTO = RateXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
RateXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
RateXSTO = RateXSTO.join(banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)
RateXSTO = RateXSTO.astype({
                    'rating 1':int,
                    'rating 2':int,
                    'rating 3':int,
                    'rating 4':int,
                    'rating 5':int,
})

# Feedback extraction
def casefold (content):
    content = re.sub("[\[\]\.\!\?\:\;\=\'\"\+\_\~\(\)\/]", '', content)
    return content
df['responses.selectedOptions'] = df['responses.selectedOptions'].apply(casefold)

replace = {"Arrive on time": "Datang tepat waktu",
           "Friendly":"Ramah",
           "Explain the cause of the disruption":"Menjelaskan penyebab gangguan",
           "The technician's appearance / uniform is presentable" : "Penampilan/ seragam teknisi rapi",
           "Resolve the disruption quickly": "Menyelesaikan gangguan dengan cepat",
           "Slow resolution of disruption":"Penyelesaikan gangguan lambat",
           "Did not arrive on time":"Tidak datang tepat waktu",
           "The technicians appearance  uniform is presentable":"Penampilan seragam teknisi rapi",
           "Explain the service properly":"Menjelaskan layanan dengan baik",
           "Resolve the installation quickly":"Menyelesaikan gangguan dengan cepat"}

for key in replace:
    df['responses.selectedOptions'] = df['responses.selectedOptions'].str.replace(key, replace[key])

feedback = df['responses.selectedOptions'].str.get_dummies(sep=',')
ordered_feedback = pd.DataFrame(feedback.sum(axis=0).sort_values(ascending=False)).index.tolist() #order column based on their value
feedback = feedback[ordered_feedback]
# FeedbackXRegion
FeedbackXRegion_raw = pd.concat([df['responses.region'], feedback], axis=1)
FeedbackXRegion = FeedbackXRegion_raw.groupby(['responses.region'], as_index=False).sum()
# FeedbackXWitel
FeedbackXWitel_raw = pd.concat([df['responses.witel'], feedback], axis=1)
FeedbackXWitel = FeedbackXWitel_raw.groupby(['responses.witel'], as_index=False).sum()
FeedbackXWitel.rename(columns = {'responses.witel': 'witel'}, inplace = True)
FeedbackXWitel = FeedbackXWitel.join(banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
# FeedbackXSTO
FeedbackXSTO_raw = pd.concat([df['responses.sto'], feedback], axis=1)
FeedbackXSTO = FeedbackXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
FeedbackXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
FeedbackXSTO = FeedbackXSTO.join(banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)

# extract date
rateXdate = pd.concat([date, rating], axis=1)
rateXdate.rename(columns = {'responses.createdAt': 'Tanggal'}, inplace = True)
rateXdate = rateXdate.groupby(['Tanggal'], as_index=False).sum()

# function to call graph
def plotlygraph(df, str_mode, str_data=None):
    if str_data=='head':
        fig = px.bar(df.head(10),
                     x=df.head(10).columns[0],
                     y=df.head(10).columns[1:-1],
                     barmode=str_mode
                     #title="Wide-Form Input",
                    )
    elif str_data=='tail':
        fig = px.bar(df.tail(10),
                     x=df.tail(10).columns[0],
                     y=df.tail(10).columns[1:-1],
                     barmode=str_mode
                    )
    else:
        fig = px.bar(df,
                     x=df.columns[0],
                     y=df.columns[1:],
                     barmode=str_mode
                    )
    st.plotly_chart(fig, use_container_width=True)   
    return


radio = st.sidebar.radio(label = '',options=('Home', 'Rating', 'Feedback'))
if radio == 'Home':
    st.sidebar.write('You selected homepage.')
    st.write(banyak_sto)
elif radio == 'Rating':
    st.header('Ringkasan Rating Teknisi Per Region')
    plotlygraph(RateXRegion, 'group')

    st.header('Ringkasan Rating Teknisi Per Witel')
    st.subheader('Top 10 Rating')
    plotlygraph(RateXWitel_fix, 'group', 'head')
    st.subheader('Bottom 10 Rating')
    plotlygraph(RateXWitel_fix, 'group', 'tail')

    st.header('Ringkasan Rating Teknisi Per STO')
    # st.write(RateXSTO)
    st.subheader('Top 10 Rating')
    plotlygraph(RateXSTO, 'group', 'head')
    st.subheader('Bottom 10 Rating')
    plotlygraph(RateXSTO, 'group', 'tail')
    
    st.header('Ringkasan Rating Teknisi Per Bulan')
    plotlygraph(rateXdate, 'group')
else:
    st.header('Ringkasan Feedback Teknisi Per Region')
    plotlygraph(FeedbackXRegion, 'group')

    st.header('Ringkasan Feedback Teknisi Per Witel')
    st.subheader('Top 10 Feedback')
    plotlygraph(FeedbackXWitel, 'group', 'head')
    st.subheader('Bottom 10 Feedback')
    plotlygraph(FeedbackXWitel, 'group', 'tail')

    st.header('Ringkasan Feedback Teknisi Per STO')
    # st.write(FeedbackXSTO)
    st.subheader('Top 10 Feedback')
    plotlygraph(FeedbackXSTO, 'group', 'head')
    st.subheader('Bottom 10 Feedback')
    plotlygraph(FeedbackXSTO, 'group', 'head')


# MAKE PPTX FILE
prs = Presentation('slidemstr.pptx')
# prs.slide_width = Inches(16)
# prs.slide_height = Inches(9)
slide_reg = prs.slide_layouts[11]

# Make fbXreg slide
def fbreg(df, str_title):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    title.text = str_title
    # title1.text_frame.paragraphs[0].font.name = "Arial"
    chart_data = ChartData()
    chart_data.categories = list(df.iloc[:,0])
    for i in range(1,len(df.columns)):
        chart_data.add_series(df.columns[i], (df.iloc[:,i]))
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = RGBColor(0,0,0)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'
    return


def addSeries(df, str_title, head=True):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    # title3.top = Inches(3)
    title.text = str_title
    chart_data = ChartData()
    if head:
        chart_data.categories = list(df.head(10).iloc[:,0])
        for i in range (1,len(df.columns)-1):
            chart_data.add_series(df.head(10).columns[i], (df.head(10).iloc[:,i]))
    else:
        chart_data.categories = list(df.tail(10).iloc[:,0])
        for i in range (1,len(df.columns)-1):
            chart_data.add_series(df.tail(10).columns[i], (df.tail(10).iloc[:,i]))

    x, y, cx, cy = Inches(1), Inches(1), Inches(12), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = 'Arial'
    return

def piechart(df):
    slide = prs.slides.add_slide(slide_reg)
    title = slide.shapes.title
    # .left .top .width .hight
    title.width.top = Inches(3), Inches(3)
    # title3.top = Inches(3)
    # title.text = str_title
    chart_data = ChartData()
    chart_data.categories = list(df.value_counts().index)
    chart_data.add_series('pie', (df.value_counts().values))
    x, y, cx, cy = Inches(0.5), Inches(4), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = chart.plots[0].data_labels
    # data_labels.show_percentage = True
    # data_labels.number_format = "0\%"
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return

slide = prs.slides.add_slide(slide_reg)


chart_data = CategoryChartData()
chart_data.categories = list(banyak_region.iloc[:,0].values)
chart_data.add_series('Regional', (banyak_region.iloc[:,1].values))

chart_data1 = CategoryChartData()
chart_data1.categories = list(banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,0].values)
chart_data1.add_series('STO', (banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,1].values))

# chart_data2 = CategoryChartData()
# chart_data2.categories = list(banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,0].values)
# chart_data1.add_series('Regional', (banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,1].values))

x, y, cx, cy = Inches(0.5), Inches(1), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart.plots[0].vary_by_categories = False

slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x=Inches(7), y=Inches(1), cx=Inches(6), cy=Inches(4.5), chart_data=chart_data1
).chart.plots[0].vary_by_categories = False

chart_data2 = ChartData()
chart_data2.categories = list(df['responses.rate'].value_counts().index)
chart_data2.add_series('pie', (df['responses.rate'].value_counts().values))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x=Inches(0.5), y=Inches(3), cx=Inches(6), cy=Inches(4.5), chart_data=chart_data2
).chart
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = chart.plots[0].data_labels
    # data_labels.show_percentage = True
    # data_labels.number_format = "0\%"
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False


# piechart(df['responses.rate'])

fbreg(RateXRegion, 'Ringkasan Rating Teknisi Per Region')
addSeries(RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Top 10")
addSeries(RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Bottom 10", False)
addSeries(RateXSTO, "Ringkasan Rating Teknisi Per STO Top 10")
addSeries(RateXSTO, "Ringkasan Rating Teknisi Per STO Bottom 10", False)
fbreg(rateXdate, 'Ringkasan Rating Teknisi Per Bulan')
fbreg(FeedbackXRegion, 'Ringkasan Feedback Teknisi Per Region')
addSeries(FeedbackXWitel, "Ringkasan Feedback Teknisi Per Witel Top 10")
addSeries(FeedbackXWitel, "Ringkasan Feedback Teknisi Per Witel Bottom 10", False)
addSeries(FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Top 10")
addSeries(FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Bottom 10", False)
# piechart(sum_avgr_ticketreop)
last_slide = prs.slides.add_slide(prs.slide_layouts[13])

output = BytesIO()
prs.save(output)
st.sidebar.header('Generate PPTX Below')
st.sidebar.download_button(
     label="Generate PPTX",
     data=output,
     file_name='technician_rate.pptx',
     mime='text/csv'
)


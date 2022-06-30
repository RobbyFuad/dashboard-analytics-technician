import tkinter as tk
import pandas as pd
import re
from tkinter import *
from tkinter import ttk
from tkinter import filedialog
from tkinter.filedialog import askopenfilename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor


class App(tk.Tk):
    
    def __init__(self):
        super().__init__()
        # configure the root window
        self.title('MyTechnician')
        self.geometry('500x150')
        p1 = PhotoImage(file = 'telkom.png')
        # Setting icon of master window
        self.iconphoto(False, p1)
        # path string
        self.path_var = tk.StringVar()
        self.path_var2 = tk.StringVar()
        # label
        self.label = ttk.Label(self, text='Powerpoint automate')
        self.label.pack()
        # button
        self.button = ttk.Button(self, text='Choose file')
        self.button['command'] = self.upload
        self.button.pack()
        
        label2 = tk.Label(self, textvariable=self.path_var, fg='red' )
        label2.pack()
        self.path_var.set("")
        
        self.button2 = ttk.Button(self, text='Generate PPTX')
        self.button2['command'] = self.ppt
        self.button2.pack()
        
        label3 = tk.Label(self, textvariable=self.path_var2, fg='red' )
        label3.pack()
        self.path_var2.set("")
        
    def upload(self):
        file = filedialog.askopenfilename(title="Select a file",
                                    filetypes=((".xlsx files", "*.xlsx"), ("all files", "*.*")))
        if file:
            self.data(file)
            self.path_var.set(file)
    
    def data(self, file):
        global df
        df = pd.read_excel(file)
        return df
    
    def preprocess(self):
        
        date = pd.to_datetime(df['responses.createdAt']).dt.strftime("%Y-%m")
        df['indiHomeNum'] = df['indiHomeNum'].astype(str)
        df['responses.region'] = 'region ' + df['indiHomeNum'].str[1]
        witel_mapper = {'JATENG UTARA  (SEMARANG)': 'SEMARANG', 'BABEL': 'BANGKA BELITUNG', 'RIDAR': 'RIAU DARATAN','KALTENG':'KALIMANTAN TENGAH',
            'JATENG SELATAN  (MAGELANG)':'MAGELANG', 'SUMUT': 'SUMATERA UTARA', 'KALBAR': 'KALIMANTAN BARAT', 'RIKEP': 'RIAU KEPULAUAN','SUMATERA SELATAN (PALEMBANG)':'PALEMBANG',
            'SUMSEL': 'SUMATERA SELATAN', 'SERANG': 'BANTEN', 'BANDUNGBRT': 'BANDUNG BARAT', 'JATENG BARAT SELATAN  (PWKT)': 'PURWOKERTO','SUMATERA BARAT (PADANG)':'PADANG',
            'KALTARA': 'KALIMANTAN UTARA', 'JATIM TENGAH  (KEDIRI)': 'KEDIRI', 'SUMBAR': 'SUMATERA BARAT', 'KALSEL': 'KALIMANTAN SELATAN','JABAR TENGAH (BANDUNG BRT)': 'BANDUNG BARAT',
            'JATENG TIMUR SELATAN  (SOLO)': 'SOLO', 'JATIM SELATAN  (MALANG)': 'MALANG', 'JATIM TENGAH TIMUR  (SIDOARJO)': 'SIDOARJO','KALTENG  (PALANGKARAYA)':'PALANGKARAYA',
            'BALI SELATAN  (DENPASAR)': 'DENPASAR', 'JAKSEL': 'JAKARTA SELATAN', 'JAKTIM': 'JAKARTA TIMUR', 'JABAR BARAT UTARA (BEKASI)': 'BEKASI','KALSEL  (BANJARMASIN)':'BANJARMASIN',
            'SULTENG': 'SULAWESI TENGAH', 'SUMUT BARAT (MEDAN)': 'MEDAN', 'BALI UTARA  (SINGARAJA)': 'SINGARAJA', 'DI YOGYAKARTA': 'YOGYAKARTA','BANTEN TIMUR (TANGERANG)': 'TANGERANG',
            'SULTRA': 'SULAWESI TENGGARA', 'NTT  (KUPANG)': 'KUPANG', 'JATIMSEL TIMUR  (PASURUAN)': 'PASURUAN', 'JABAR BARAT (BOGOR)': 'BOGOR','JAKUT':'JAKARTA UTARA',
            'DATEL MAKASAR':'MAKASAR','JAKBAR': 'JAKARTA BARAT','JATENG BARAT UTARA  (PKLG)':'PEKALONGAN','JATENG TIMUR UTARA  (KUDUS)':'KUDUS',
            'JATIM BARAT  (MADIUN)':'MADIUN','KALTIMSEL  (BALIKPAPAN)':'BALIKPAPAN','NTB  (MATARAM)':'MATARAM',
            'KALTIMTENG  (SAMARINDA)':'SAMARINDA', 'KALTIMUT  (TARAKAN)':'TARAKAN', 'DATEL PONTIANAK':'PONTIANAK', 'JABAR TIMUR (CIREBON)':'CIREBON', 'JABAR UTARA (KARAWANG)':'KARAWANG',
            'DATEL JAKARTA TIMUR':'JAKARTA TIMUR', 'BANGKA BELITUNG (P.PINANG)':'P.PINANG', 'RIAU KEPULAUAN (BATAM)':'BATAM', 'SUMUT TIMUR (PEMATANG SIANTAR)':'PEMATANG SIANTAR'}

        df['responses.witel'] = df['responses.witel'].map(witel_mapper).fillna(df['responses.witel'])
        witel = df['responses.witel'].unique()


        self.banyak_witel = df['responses.witel'].value_counts().rename_axis('witel').reset_index(name='count').sort_values(by=['witel'])
        self.banyak_sto = df['responses.sto'].value_counts().rename_axis('sto').reset_index(name='count').sort_values(by=['sto'])
        self.banyak_region = df['responses.region'].value_counts().rename_axis('region').reset_index(name='count').sort_values(by=['region'])

        # Rating Extraction
        rating = pd.get_dummies(df['responses.rate'], columns=['responses.rate'], prefix='rating', prefix_sep=' ')

        # average per tiket reopen
        avgr_ticketreop = pd.concat([df['responses.reopenCount'], rating], axis=1)
        sum_avgr_ticketreop = avgr_ticketreop.groupby(['responses.reopenCount'], as_index=False).sum()
        mean_avgr_ticketreop = avgr_ticketreop.groupby(['responses.reopenCount'], as_index=False).mean()
        # RateXRegion
        RateXRegion_raw = pd.concat([df['responses.region'], rating], axis=1)
        self.RateXRegion = RateXRegion_raw.groupby(['responses.region'], as_index=False).sum()
        # RateXWitel
        RateXWitel_raw = pd.concat([df['responses.witel'], rating], axis=1)
        RateXWitel = RateXWitel_raw.groupby(['responses.witel'], as_index=False).sum()
        RateXWitel.rename(columns = {'responses.witel': 'witel'}, inplace = True)
        self.RateXWitel_fix = RateXWitel.join(self.banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
        # RateXSTO
        RateXSTO_raw = pd.concat([df['responses.sto'], rating], axis=1)
        RateXSTO = RateXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
        RateXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
        RateXSTO = RateXSTO.join(self.banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)
        self.RateXSTO = RateXSTO.astype({
                    'rating 1':int,
                    'rating 2':int,
                    'rating 3':int,
                    'rating 4':int,
                    'rating 5':int,
        })

        # Feedback extraction
        def casefold (content):
            content = re.sub("[\[\]\.\!\?\:\;\=\'\"\+\_\~\(\)\/]", '', content)
            return content
        df['responses.selectedOptions'] = df['responses.selectedOptions'].apply(casefold)

        replace = {"Arrive on time": "Datang tepat waktu",
                   "Friendly":"Ramah",
                   "Explain the cause of the disruption":"Menjelaskan penyebab gangguan",
                   "The technician's appearance / uniform is presentable" : "Penampilan/ seragam teknisi rapi",
                   "Resolve the disruption quickly": "Menyelesaikan gangguan dengan cepat",
                   "Slow resolution of disruption":"Penyelesaikan gangguan lambat",
                   "Did not arrive on time":"Tidak datang tepat waktu",
                   "The technicians appearance  uniform is presentable":"Penampilan seragam teknisi rapi",
                   "Explain the service properly":"Menjelaskan layanan dengan baik",
                   "Resolve the installation quickly":"Menyelesaikan gangguan dengan cepat",
                   "Menyelesaikan ganggaun dengan cepat":"Menyelesaikan gangguan dengan cepat"}

        for key in replace:
            df['responses.selectedOptions'] = df['responses.selectedOptions'].str.replace(key, replace[key])

        feedback = df['responses.selectedOptions'].str.get_dummies(sep=',')
        ordered_feedback = pd.DataFrame(feedback.sum(axis=0).sort_values(ascending=False)).index.tolist() #order column based on their value
        feedback = feedback[ordered_feedback]
        # FeedbackXRegion
        FeedbackXRegion_raw = pd.concat([df['responses.region'], feedback], axis=1)
        self.FeedbackXRegion = FeedbackXRegion_raw.groupby(['responses.region'], as_index=False).sum()
        # FeedbackXWitel
        FeedbackXWitel_raw = pd.concat([df['responses.witel'], feedback], axis=1)
        FeedbackXWitel = FeedbackXWitel_raw.groupby(['responses.witel'], as_index=False).sum()
        FeedbackXWitel.rename(columns = {'responses.witel': 'witel'}, inplace = True)
        self.FeedbackXWitel = FeedbackXWitel.join(self.banyak_witel.set_index('witel'), on='witel').sort_values(by=['count'], ascending=False)
        # FeedbackXSTO
        FeedbackXSTO_raw = pd.concat([df['responses.sto'], feedback], axis=1)
        FeedbackXSTO = FeedbackXSTO_raw.groupby(['responses.sto'], as_index=False).sum()
        FeedbackXSTO.rename(columns = {'responses.sto': 'sto'}, inplace = True)
        self.FeedbackXSTO = FeedbackXSTO.join(self.banyak_sto.set_index('sto'), on='sto').sort_values(by=['count'], ascending=False)
        # extract date
        rateXdate = pd.concat([date, rating], axis=1)
        rateXdate.rename(columns = {'responses.createdAt': 'Tanggal'}, inplace = True)
        self.rateXdate = rateXdate.groupby(['Tanggal'], as_index=False).sum()        
        # FeedbackXSTO
        FeedbackXrate_raw = pd.concat([df['responses.rate'], feedback], axis=1)
        FeedbackXrate_raw['responses.rate'] = 'rating ' + FeedbackXrate_raw['responses.rate'].apply(str)
        self.FeedbackXrate = FeedbackXrate_raw.groupby(['responses.rate'], as_index=False).sum()
        self.FeedbackXrate.rename(columns = {'responses.rate': 'rating'}, inplace = True)
        
        
        
    def ppt(self):
        self.preprocess()
        self.prs = Presentation('slidemstr.pptx')
        slide_reg = self.prs.slide_layouts[11]
        slide = self.prs.slides.add_slide(slide_reg)
        
        # Make fbXreg slide
        def fbreg(df, str_title):
            slide = self.prs.slides.add_slide(slide_reg)
            title = slide.shapes.title
            # .left .top .width .hight
            title.width.top = Inches(3), Inches(3)
            title.text = str_title
            # title1.text_frame.paragraphs[0].font.name = "Arial"
            chart_data = ChartData()
            chart_data.categories = list(df.iloc[:,0])
            for i in range(1,len(df.columns)):
                chart_data.add_series(df.columns[i], (df.iloc[:,i]))
            
            x, y, cx, cy = Inches(0.3), Inches(1), Inches(12.7), Inches(6)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            
            plot = chart.plots[0]
            plot.overlap = -100
            plot.gap_width = 100
            plot.has_data_labels = True
            
            data_labels = plot.data_labels
            data_labels.font.size = Pt(7)
            data_labels.font.color.rgb = RGBColor(0,0,0)
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = 'Arial'
            return


        def addSeries(df, str_title, head=True):
            slide = self.prs.slides.add_slide(slide_reg)
            title = slide.shapes.title
            # .left .top .width .hight
            title.width.top = Inches(3), Inches(3)
            # title3.top = Inches(3)
            title.text = str_title
            chart_data = ChartData()
            if head:
                chart_data.categories = list(df.head(10).iloc[:,0])
                for i in range (1,len(df.columns)-1):
                    chart_data.add_series(df.head(10).columns[i], (df.head(10).iloc[:,i]))
            else:
                chart_data.categories = list(df.tail(10).iloc[:,0])
                for i in range (1,len(df.columns)-1):
                    chart_data.add_series(df.tail(10).columns[i], (df.tail(10).iloc[:,i]))

            x, y, cx, cy = Inches(0.3), Inches(1), Inches(12.7), Inches(6)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            plot = chart.plots[0]
            plot.overlap = -100
            plot.gap_width = 100
            plot.has_data_labels = True
            
            data_labels = plot.data_labels
            data_labels.font.size = Pt(7)
            data_labels.font.color.rgb = RGBColor(0,0,0)
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.minor_unit = 1
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = 'Arial'
            return
        
        
        title = slide.shapes.title
        title.width.top = Inches(3), Inches(3)
        title.text = 'Ringkasan Teknisi'
        
        chart_data = CategoryChartData()
        chart_data.categories = list(self.banyak_region.iloc[:,0].values)
        chart_data.add_series('Regional', (self.banyak_region.iloc[:,1].values))

        chart_data1 = CategoryChartData()
        chart_data1.categories = list(self.banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,0].values)
        chart_data1.add_series('Top 10 STO', (self.banyak_sto.sort_values(by=['count'], ascending=False).head(10).iloc[:,1].values))

        chart_data3 = CategoryChartData()
        chart_data3.categories = list(self.banyak_witel.sort_values(by=['count'], ascending=False).head(10).iloc[:,0].values)
        chart_data3.add_series('Top 10 Witel', (self.banyak_witel.sort_values(by=['count'], ascending=False).head(10).iloc[:,1].values))


        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x=Inches(0.3), y=Inches(0.8), cx=Inches(6), cy=Inches(3.5), chart_data=chart_data
        ).chart
        plot = chart.plots[0]
        plot.vary_by_categories = False
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.value_axis.tick_labels.font.size = Pt(11)
        # chart_title.text_frame.text = 'yourtitle'
        # chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12) 
        # this is font size title but seems to overwrite the title defined before, need to define title again

        chart1 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x=Inches(7), y=Inches(0.8), cx=Inches(6), cy=Inches(3.5), chart_data=chart_data1
        ).chart
        plot = chart1.plots[0]
        plot.vary_by_categories = False
        chart1.category_axis.tick_labels.font.size = Pt(11)
        chart1.value_axis.tick_labels.font.size = Pt(11)

        chart3 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x=Inches(5), y=Inches(4.2), cx=Inches(8), cy=Inches(3.5), chart_data=chart_data3
        ).chart
        plot = chart3.plots[0]
        plot.vary_by_categories = False
        chart3.category_axis.tick_labels.font.size = Pt(9)
        chart3.value_axis.tick_labels.font.size = Pt(11)

        chart_data2 = ChartData()
        chart_data2.categories = list(df['responses.rate'].value_counts().index)
        chart_data2.add_series('Persentase Rating', (df['responses.rate'].value_counts().values))
        chart2 = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x=Inches(0.3), y=Inches(4.2), cx=Inches(4), cy=Inches(3), chart_data=chart_data2
        ).chart
        plot = chart2.plots[0]
        plot.has_data_labels = True
        data_labels = chart2.plots[0].data_labels
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(9)
        data_labels.show_percentage = True
        data_labels.number_format = '0.00%'
        chart2.has_legend = True
        chart2.legend.font.size = Pt(10)
        chart2.legend.position = XL_LEGEND_POSITION.RIGHT
        chart2.legend.include_in_layout = True

        fbreg(self.FeedbackXrate, 'Ringkasan Feedback Teknisi Per Rating')
        fbreg(self.RateXRegion, 'Ringkasan Rating Teknisi Per Region')
        addSeries(self.RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Top 10")
        addSeries(self.RateXWitel_fix, "Ringkasan Rating Teknisi Per Witel Bottom 10", False)
        addSeries(self.RateXSTO, "Ringkasan Rating Teknisi Per STO Top 10")
        addSeries(self.RateXSTO, "Ringkasan Rating Teknisi Per STO Bottom 10", False)
        fbreg(self.rateXdate, 'Ringkasan Rating Teknisi Per Bulan')
        fbreg(self.FeedbackXRegion, 'Ringkasan Feedback Teknisi Per Region')
        addSeries(self.FeedbackXWitel, "Ringkasan Feedback Teknisi Per Witel Top 10")
        addSeries(self.FeedbackXWitel, "Ringkasan Feedback Teknisi Per Witel Bottom 10", False)
        addSeries(self.FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Top 10")
        addSeries(self.FeedbackXSTO, "Ringkasan Feedback Teknisi Per STO Bottom 10", False)
        # piechart(sum_avgr_ticketreop)
        last_slide = self.prs.slides.add_slide(self.prs.slide_layouts[13])
        
        self.prs.save('technician_rate.pptx')
        self.path_var2.set("PPTX Updated")


if __name__ == "__main__":
  app = App()
  app.mainloop()